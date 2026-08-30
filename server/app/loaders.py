import hashlib
import json
from bisect import bisect_right
from pathlib import Path
from typing import Callable, Optional

TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv", ".ini",
    ".cfg", ".toml", ".yaml", ".yml", ".json", ".xml", ".svg",
    ".py", ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx", ".java", ".kt",
    ".c", ".h", ".cpp", ".hpp", ".cs", ".go", ".rs", ".rb", ".php",
    ".swift", ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd", ".sql",
    ".ipynb", ".env.example", ".gitignore", ".dockerfile", ".makefile",
}
DOC_EXTS = {".pdf", ".docx", ".pptx"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}

SUPPORTED_EXTS = TEXT_EXTS | DOC_EXTS | IMAGE_EXTS
SUPPORTED_NAMES = {".env.example", ".gitignore", "dockerfile", "makefile"}

IGNORED_DIRS = {
    ".git", ".svn", ".hg", "node_modules", "__pycache__", ".venv",
    "venv", ".env", "dist", "build", ".next", ".nuxt", "target",
    "bin", "obj", "coverage", ".idea", ".vscode", ".gradle",
    ".pytest_cache", ".mypy_cache", ".ruff_cache", "chroma",
}

MAX_FILE_BYTES = 25 * 1024 * 1024

CHUNKER_VERSION = "v2"

def file_signature(path: Path) -> str:
    stat = path.stat()
    raw = f"{CHUNKER_VERSION}|{path}|{stat.st_size}|{stat.st_mtime_ns}"
    return hashlib.sha1(raw.encode()).hexdigest()


def is_supported_file(path: Path) -> bool:
    return path.name.lower() in SUPPORTED_NAMES or path.suffix.lower() in SUPPORTED_EXTS


def chunk_id(path: str, idx: int) -> str:
    return hashlib.sha1(path.encode()).hexdigest()[:16] + f"-{idx}"


def chunk_lines(
    lines: list[str], max_chars: int = 1200, overlap_lines: int = 3
) -> list[tuple[int, int, str]]:
    chunks: list[tuple[int, int, str]] = []
    buf: list[str] = []
    start = 0

    def flush():
        nonlocal buf, start
        if not buf:
            return
        text = "\n".join(buf).strip()
        if text:
            end = start + len(buf) - 1
            chunks.append((start + 1, end + 1, text))
            keep: list[str] = []
            budget = 0
            for ln in reversed(buf):
                if budget + len(ln) > 200 or len(keep) >= overlap_lines:
                    break
                keep.insert(0, ln)
                budget += len(ln)
            start = end - len(keep) + 1
            buf = list(keep)
        else:
            buf = []

    for i, line in enumerate(lines):
        if sum(len(l) for l in buf) + len(line) > max_chars and buf:
            flush()
        if not buf:
            start = i
        buf.append(line)
    while buf:
        text = "\n".join(buf).strip()
        if text:
            chunks.append((start + 1, start + len(buf), text))
        break
    return chunks


def _read_text_lines(path: Path) -> Optional[list[str]]:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
        if text.count("\x00") > 5:
            return None
        return text.splitlines()
    except Exception:
        return None


def _read_ipynb(path: Path) -> Optional[list[str]]:
    try:
        nb = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    except Exception:
        return None
    lines: list[str] = []
    for n, cell in enumerate(nb.get("cells", [])):
        src = "".join(cell.get("source", []))
        kind = cell.get("cell_type", "code")
        lines.append(f"##### {kind} cell {n + 1} #####")
        lines.extend(src.splitlines())
    return lines


def extract_pdf(path: Path) -> list[dict]:
    import fitz

    out = []
    with fitz.open(str(path)) as doc:
        for pno in range(doc.page_count):
            page = doc[pno]
            text = page.get_text("text")
            for s, e, t in chunk_lines(text.splitlines()):
                out.append({"text": t, "page": pno + 1,
                            "line_start": s, "line_end": e})
    return out


def extract_docx(path: Path) -> list[dict]:
    from docx import Document

    doc = Document(str(path))
    lines = [p.text for p in doc.paragraphs]
    return [
        {"text": t, "page": None, "line_start": s, "line_end": e}
        for s, e, t in chunk_lines(lines)
    ]


def extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for sno, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs)
                    if txt.strip():
                        parts.append(txt)
        for s, e, t in chunk_lines(parts):
            out.append({"text": t, "page": sno + 1,
                        "line_start": s, "line_end": e})
    return out


def extract_file(path: Path) -> Optional[list[dict]]:
    ext = path.suffix.lower()
    name = path.name.lower()
    try:
        if ext == ".pdf":
            return extract_pdf(path)
        if ext == ".docx":
            return extract_docx(path)
        if ext == ".pptx":
            return extract_pptx(path)
        if ext == ".ipynb":
            lines = _read_ipynb(path) or []
        elif ext in TEXT_EXTS or name in SUPPORTED_NAMES:
            got = _read_text_lines(path)
            if got is None:
                return None
            lines = got
        else:
            return None
        return [
            {"text": t, "page": None, "line_start": s, "line_end": e}
            for s, e, t in chunk_lines(lines)
        ]
    except Exception:
        return None


def build_chunks(path: Path, root: Path) -> dict:
    sig = file_signature(path)
    rel = str(path.relative_to(root)).replace("\\", "/")
    pieces = extract_file(path)
    if not pieces:
        return {}
    chunks = []
    for i, piece in enumerate(pieces):
        meta = {
            "path": str(path),
            "rel_path": rel,
            "root": str(root),
            "name": path.name,
            "ext": path.suffix.lower().lstrip(".") or "txt",
            "line_start": piece["line_start"],
            "line_end": piece["line_end"],
            "page": piece["page"] if piece["page"] is not None else -1,
            "signature": sig,
            "chunker_version": CHUNKER_VERSION,
        }
        chunks.append({
            "id": chunk_id(str(path), i),
            "text": piece["text"],
            "meta": meta,
        })
    return {"signature": sig, "chunks": chunks}


def describe_image_chunk(path: Path, root: Path, describe_fn: Callable[[bytes, str, str], str]) -> Optional[dict]:
    sig = file_signature(path)
    mime_map = {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".webp": "image/webp", ".gif": "image/gif", ".bmp": "image/bmp",
    }
    mime = mime_map.get(path.suffix.lower())
    if not mime:
        return None
    try:
        data = path.read_bytes()
        desc = describe_fn(data, mime, str(path.relative_to(root)))
    except Exception:
        return None
    rel = str(path.relative_to(root)).replace("\\", "/")
    return {
        "id": chunk_id(str(path), 0),
        "text": desc,
        "meta": {
            "path": str(path), "rel_path": rel, "root": str(root),
            "name": path.name, "ext": path.suffix.lower().lstrip("."),
            "line_start": 0, "line_end": 0, "page": -1,
            "signature": sig, "type": "image",
        },
    }


def snippet_at(text: str, line_start: int) -> str:
    return text[:400]
